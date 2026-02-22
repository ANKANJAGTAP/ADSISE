"""
Generate a 7-slide PowerPoint presentation for the Expense Tracker PWA project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette (teal / dark theme inspired by the app) ──
TEAL        = RGBColor(0x00, 0x96, 0x88)   # primary
TEAL_DARK   = RGBColor(0x00, 0x79, 0x6B)
TEAL_LIGHT  = RGBColor(0xB2, 0xDF, 0xDB)
DARK_BG     = RGBColor(0x1E, 0x1E, 0x2E)   # slide background
DARK_BG2    = RGBColor(0x2A, 0x2A, 0x3C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
GOLD        = RGBColor(0xFF, 0xD5, 0x4F)
ORANGE      = RGBColor(0xFF, 0x8A, 0x65)
CARD_BG     = RGBColor(0x2D, 0x2D, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ────────────────────── helper functions ──────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=LIGHT_GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=TEAL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_card(slide, left, top, w, h, title, body_lines,
             accent=TEAL, title_size=18, body_size=14):
    """Card with accent top border."""
    # accent bar
    add_rect(slide, left, top, w, Pt(5), accent)
    # card bg
    add_shape(slide, left, top + Pt(5), w, h - Pt(5), CARD_BG)
    # title
    add_text_box(slide, left + Inches(0.2), top + Pt(14),
                 w - Inches(0.4), Inches(0.4),
                 title, title_size, WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    # body
    y = top + Pt(14) + Inches(0.45)
    add_bullet_list(slide, left + Inches(0.2), y,
                    w - Inches(0.4), h - Inches(1),
                    body_lines, body_size, LIGHT_GRAY, spacing=Pt(4))


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# decorative circles
add_circle(slide, Inches(-1), Inches(-1), Inches(4), TEAL_DARK)
add_circle(slide, Inches(10.5), Inches(5), Inches(4), TEAL_DARK)

# small accent line
add_accent_line(slide, Inches(4.2), Inches(2.3), Inches(4.9), TEAL)

# Title
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.4),
             "💰  Expense Tracker PWA", font_size=48, color=WHITE,
             bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Subtitle
add_text_box(slide, Inches(2), Inches(4), Inches(9), Inches(0.8),
             "A Progressive Web Application to Track Your Daily Expenses — Online & Offline",
             font_size=22, color=LIGHT_GRAY, bold=False,
             alignment=PP_ALIGN.CENTER)

# Tech badges row
badges = ["React 19", "Vite 7", "IndexedDB", "Service Worker", "PWA"]
badge_w = Inches(1.8)
total_badges_w = len(badges) * badge_w + (len(badges)-1) * Inches(0.2)
start_x = (SLIDE_W - total_badges_w) // 2
for i, badge in enumerate(badges):
    x = start_x + i * (badge_w + Inches(0.2))
    add_shape(slide, x, Inches(5.2), badge_w, Inches(0.5), TEAL_DARK, TEAL)
    add_text_box(slide, x, Inches(5.22), badge_w, Inches(0.5),
                 badge, 14, TEAL_LIGHT, True, PP_ALIGN.CENTER)

# Student info card at bottom-left
add_shape(slide, Inches(0.5), Inches(6.0), Inches(5.5), Inches(1.2), CARD_BG, TEAL_DARK)
info_lines = [
    "Name: Ankan Jagtap    |    PRN: 23510058",
    "Topic No: 13    |    Topic: Progressive Web Application (PWA) Technology",
]
add_bullet_list(slide, Inches(0.7), Inches(6.1), Inches(5.2), Inches(1.0),
                info_lines, font_size=13, color=TEAL_LIGHT, spacing=Pt(4))

# Footer
add_text_box(slide, Inches(6.5), Inches(6.6), Inches(6.5), Inches(0.4),
             "Progressive Web Application Technology  •  February 2026",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION / PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(3), TEAL)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Introduction & Problem Statement", 36, WHITE, True)

# Left column – Problem
add_card(slide,
         Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
         "🔍  The Problem",
         [
             "• People struggle to keep track of daily expenses",
             "• Traditional apps require constant internet",
             "• Data loss risk when offline or switching devices",
             "• No lightweight solution that works everywhere",
         ],
         accent=ORANGE, body_size=15)

# Right column – Solution
add_card(slide,
         Inches(7), Inches(2.2), Inches(5.5), Inches(4.5),
         "✅  Our Solution",
         [
             "• A PWA that works fully offline",
             "• Data persists in browser via IndexedDB",
             "• Installable on any device — no app store needed",
             "• Fast, lightweight, and responsive UI",
             "• Real-time online/offline status indicator",
         ],
         accent=TEAL, body_size=15)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(3), TEAL)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Technology Stack", 36, WHITE, True)

techs = [
    ("⚛️  React 19",      "Component-based UI with hooks\n(useState, useEffect)",            TEAL),
    ("⚡  Vite 7",         "Lightning-fast build tool\nwith HMR for development",             RGBColor(0x64, 0x6C, 0xFF)),
    ("🗄️  IndexedDB",     "Client-side NoSQL database\nfor offline data persistence",        ORANGE),
    ("⚙️  Service Worker", "Caches assets & enables\noffline-first experience",               RGBColor(0x66, 0xBB, 0x6A)),
    ("📱  PWA Manifest",   "Makes app installable with\ncustom icons & splash screen",        GOLD),
    ("📦  idb Library",    "Promise-based wrapper for\ncleaner IndexedDB async/await",        RGBColor(0xAB, 0x47, 0xBC)),
]

cols, rows = 3, 2
card_w, card_h = Inches(3.6), Inches(2.0)
gap_x, gap_y = Inches(0.4), Inches(0.35)
start_x = Inches(0.8)
start_y = Inches(2.1)

for idx, (title, desc, accent) in enumerate(techs):
    r, c = divmod(idx, cols)
    x = start_x + c * (card_w + gap_x)
    y = start_y + r * (card_h + gap_y)
    # card
    add_rect(slide, x, y, card_w, Pt(5), accent)
    add_shape(slide, x, y + Pt(5), card_w, card_h - Pt(5), CARD_BG)
    add_text_box(slide, x + Inches(0.2), y + Pt(14),
                 card_w - Inches(0.4), Inches(0.35),
                 title, 17, WHITE, True)
    add_text_box(slide, x + Inches(0.2), y + Pt(14) + Inches(0.4),
                 card_w - Inches(0.4), Inches(1.0),
                 desc, 14, LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE & WORKFLOW
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(3), TEAL)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Architecture & Workflow", 36, WHITE, True)

# Flow boxes
flow_items = [
    ("User Interface",     "React Components\n(App, ExpenseForm,\nExpenseList)",   TEAL),
    ("State Management",   "React Hooks\n(useState, useEffect)\nOnline/Offline",   RGBColor(0x64, 0x6C, 0xFF)),
    ("Data Layer",         "IndexedDB via idb\nCRUD Operations\n(Add, Read, Delete)", ORANGE),
    ("Caching Layer",      "Service Worker\nCache-first (static)\nNetwork-first (dynamic)", RGBColor(0x66, 0xBB, 0x6A)),
    ("PWA Shell",          "Web App Manifest\nInstallable App\nOffline Support",    GOLD),
]

box_w = Inches(2.1)
box_h = Inches(2.8)
gap = Inches(0.25)
total_w = len(flow_items) * box_w + (len(flow_items)-1) * gap
sx = (SLIDE_W - total_w) // 2
sy = Inches(2.3)

for i, (title, desc, accent) in enumerate(flow_items):
    x = sx + i * (box_w + gap)
    # card
    add_rect(slide, x, sy, box_w, Pt(5), accent)
    add_shape(slide, x, sy + Pt(5), box_w, box_h - Pt(5), CARD_BG)
    # number circle
    add_circle(slide, x + (box_w - Inches(0.45))//2, sy + Pt(16), Inches(0.45), accent)
    add_text_box(slide, x + (box_w - Inches(0.45))//2, sy + Pt(18),
                 Inches(0.45), Inches(0.4),
                 str(i+1), 16, DARK_BG, True, PP_ALIGN.CENTER)
    # title
    add_text_box(slide, x + Inches(0.1), sy + Pt(16) + Inches(0.55),
                 box_w - Inches(0.2), Inches(0.4),
                 title, 15, WHITE, True, PP_ALIGN.CENTER)
    # desc
    add_text_box(slide, x + Inches(0.1), sy + Pt(16) + Inches(1.0),
                 box_w - Inches(0.2), Inches(1.3),
                 desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    # arrow between boxes
    if i < len(flow_items) - 1:
        arrow_x = x + box_w + Inches(0.02)
        add_text_box(slide, arrow_x, sy + box_h//2 - Inches(0.15),
                     Inches(0.22), Inches(0.35),
                     "➤", 18, TEAL, False, PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "All data flows are local — no backend server required. The app is fully client-side.",
             14, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — KEY FEATURES
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(3), TEAL)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Key Features", 36, WHITE, True)

features = [
    ("🌐", "Offline-First",       "Works without internet.\nData stored in IndexedDB."),
    ("📱", "Installable PWA",     "Add to home screen on\nany device — no app store."),
    ("➕", "Add Expenses",        "Quick form with title &\namount validation."),
    ("🗑️", "Delete Expenses",    "Remove entries with a\nsingle click."),
    ("💰", "Total Calculation",   "Auto-computed total of\nall expenses in real-time."),
    ("🟢", "Online/Offline Badge","Live status indicator\nshowing connectivity state."),
    ("⚡", "Blazing Fast",        "Vite + cached assets =\ninstant load times."),
    ("🔒", "Local & Private",     "No server — your data\nnever leaves the browser."),
]

cols, rows = 4, 2
fw, fh = Inches(2.8), Inches(2.0)
gx, gy = Inches(0.3), Inches(0.3)
sx = Inches(0.8)
sy2 = Inches(2.1)

for idx, (icon, title, desc) in enumerate(features):
    r, c = divmod(idx, cols)
    x = sx + c * (fw + gx)
    y = sy2 + r * (fh + gy)
    add_shape(slide, x, y, fw, fh, CARD_BG, TEAL_DARK)
    # icon
    add_text_box(slide, x, y + Inches(0.15), fw, Inches(0.45),
                 icon, 28, WHITE, False, PP_ALIGN.CENTER)
    # title
    add_text_box(slide, x + Inches(0.1), y + Inches(0.6),
                 fw - Inches(0.2), Inches(0.35),
                 title, 15, TEAL_LIGHT, True, PP_ALIGN.CENTER)
    # desc
    add_text_box(slide, x + Inches(0.1), y + Inches(1.0),
                 fw - Inches(0.2), Inches(0.8),
                 desc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — PROJECT STRUCTURE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(3), TEAL)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             "Project Structure & Code Overview", 36, WHITE, True)

# Left – file tree
tree_text = (
    "expense-tracker-pwa/\n"
    "├── public/\n"
    "│   ├── manifest.json      ← PWA config\n"
    "│   ├── sw.js              ← Service Worker\n"
    "│   └── icon-*.png         ← App icons\n"
    "├── src/\n"
    "│   ├── main.jsx           ← Entry point\n"
    "│   ├── App.jsx            ← Root component\n"
    "│   ├── swRegistration.js  ← SW registration\n"
    "│   ├── components/\n"
    "│   │   ├── ExpenseForm.jsx\n"
    "│   │   └── ExpenseList.jsx\n"
    "│   └── db/\n"
    "│       └── db.js          ← IndexedDB CRUD\n"
    "├── index.html\n"
    "├── vite.config.js\n"
    "└── package.json"
)
add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.3), Inches(4.8), CARD_BG, TEAL_DARK)
add_text_box(slide, Inches(1.0), Inches(2.1), Inches(5), Inches(0.35),
             "📂  File Tree", 16, TEAL_LIGHT, True)
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(5), Inches(4.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = tree_text
p.font.size = Pt(12)
p.font.color.rgb = LIGHT_GRAY
p.font.name = "Consolas"

# Right – component descriptions
right_cards = [
    ("App.jsx",          "Root component • manages state\n• loads expenses from IndexedDB\n• tracks online/offline status"),
    ("ExpenseForm.jsx",  "Controlled form • validates\ntitle & amount • delegates\npersistence to parent via props"),
    ("ExpenseList.jsx",  "Renders expense list • computes\ntotal • delete button per item\n• date formatting (en-IN)"),
    ("db.js",            "IndexedDB CRUD via idb library\n• addExpense • getAllExpenses\n• deleteExpense"),
]

ry = Inches(2.0)
for title, desc in right_cards:
    add_shape(slide, Inches(6.6), ry, Inches(5.8), Inches(1.05), CARD_BG, TEAL_DARK)
    add_text_box(slide, Inches(6.8), ry + Inches(0.05),
                 Inches(2), Inches(0.3),
                 title, 14, TEAL_LIGHT, True)
    add_text_box(slide, Inches(6.8), ry + Inches(0.3),
                 Inches(5.4), Inches(0.7),
                 desc, 12, LIGHT_GRAY)
    ry += Inches(1.2)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — THANK YOU
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# Large decorative circles
add_circle(slide, Inches(0.5), Inches(0.5), Inches(3), TEAL_DARK)
add_circle(slide, Inches(9.5), Inches(4.5), Inches(3.5), TEAL_DARK)
add_circle(slide, Inches(-0.5), Inches(5.5), Inches(2), RGBColor(0x00, 0x5F, 0x56))
add_circle(slide, Inches(11.5), Inches(-0.5), Inches(2), RGBColor(0x00, 0x5F, 0x56))

# Center content
add_text_box(slide, Inches(2), Inches(1.8), Inches(9), Inches(1.2),
             "🙏", font_size=60, color=WHITE,
             bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.0), Inches(9), Inches(1.0),
             "Thank You!", font_size=54, color=GOLD,
             bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_accent_line(slide, Inches(5), Inches(4.1), Inches(3.3), TEAL)

add_text_box(slide, Inches(2), Inches(4.4), Inches(9), Inches(0.6),
             "Expense Tracker PWA — Built with React, Vite & IndexedDB",
             font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
             "Progressive Web Application Technology  •  February 2026",
             font_size=16, color=TEAL_LIGHT, alignment=PP_ALIGN.CENTER)

# Decorative bottom badge
add_shape(slide, Inches(5.2), Inches(5.8), Inches(2.9), Inches(0.5), TEAL_DARK, TEAL)
add_text_box(slide, Inches(5.2), Inches(5.82), Inches(2.9), Inches(0.5),
             "Questions?  Let's discuss!", 15, WHITE, True, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = "/home/ankan/Progressive Web Application (PWA) Technology/expense-tracker-pwa/Expense_Tracker_PWA_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
